from typing import List, Dict
import os
import json
import shutil
import tempfile
from pptx import Presentation
from pptx.util import Pt
import xml.etree.ElementTree as ET

from .file_manager import FileManager
from .xml_handler import XMLHandler

import logging
import re

CHINESE_PATTERN = re.compile(r'[\u4e00-\u9fff\u3400-\u4dbf\u20000-\u2a6df]')


def has_chinese(s: str) -> bool:
    return bool(CHINESE_PATTERN.search(s))


def is_chinese_char(char):
    if len(char) != 1:
        return False
    code_point = ord(char)
    return (
            0x4E00 <= code_point <= 0x9FFF or
            0x3400 <= code_point <= 0x4DBF or
            0x20000 <= code_point <= 0x2A6DF
    )


def init_logger():
    log_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'logs')
    os.makedirs(log_dir, exist_ok=True)

    log_file = os.path.join(log_dir, 'ppt_processor.log')

    logger = logging.getLogger('PPTProcessor')
    logger.setLevel(logging.INFO)

    formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')

    from logging.handlers import RotatingFileHandler
    file_handler = RotatingFileHandler(
        log_file,
        maxBytes=10 * 1024 * 1024,
        backupCount=3,
        encoding='utf-8'
    )
    file_handler.setFormatter(formatter)
    logger.addHandler(file_handler)

    return logger

logger = init_logger()


class PPTProcessor:
    def __init__(self, status_callback=None):
        self.file_manager = FileManager()
        self.xml_handler = XMLHandler()
        self.font_configs = self.load_font_config()
        self.status_callback = status_callback

    def load_font_config(self):
        config_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'font_config.json')
        try:
            with open(config_path, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception as e:
            logger.error(f"加载字体配置失败: {str(e)}")
            return {}

    def process_ppt(self, input_path: str, output_path: str, gradient_config: List[Dict], font_size: str = None,
                    font_name: str = None):
        try:
            temp_dir = tempfile.mkdtemp()
            temp_ppt = os.path.join(temp_dir, "temp_processed.pptx")

            try:
                self.status_callback("开始字体和字号替换...")
                self._process_font_replacement(input_path, temp_ppt, self.font_configs)
                self.status_callback("开始应用渐变效果...")
                self._process_gradient_effects(temp_ppt, gradient_config, font_size, font_name)
                shutil.copy2(temp_ppt, output_path)
                self.status_callback(f"处理完成: {output_path}")

                return True

            finally:
                if os.path.exists(temp_dir):
                    shutil.rmtree(temp_dir)

        except Exception as e:
            logger.error(f"处理PPT时出错: {str(e)}")
            import traceback
            traceback.print_exc()
            return False

    def _process_font_replacement(self, input_path: str, output_path: str, font_configs: Dict):
        prs = Presentation(input_path)

        if not font_configs:
            logger.warning("没有有效的字体配置，跳过字体替换")
            prs.save(output_path)
            return

        for slide_idx, slide in enumerate(prs.slides):
            logger.info(f"处理幻灯片 {slide_idx + 1}")

            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        self._process_text_frame(shape.text_frame, font_configs)
                    elif hasattr(shape, "table") and shape.table:
                        self._process_table(shape.table, font_configs)
                except Exception as e:
                    logger.error(f"处理形状时出错: {str(e)}")
                    continue

        try:
            prs.save(output_path)
            logging.info(f"字体替换完成，已保存到: {output_path}")
        except Exception as e:
            logger.error(f"保存文件时出错: {str(e)}")
            shutil.copy2(input_path, output_path)

    def _process_text_frame(self, text_frame, font_configs: Dict):
        try:
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip():
                        self._apply_font_config_to_run(run, font_configs)
        except Exception as e:
            logger.error(f"处理文本框架时出错: {str(e)}")

    def _process_table(self, table, font_configs: Dict):
        try:
            for row in table.rows:
                for cell in row.cells:
                    if hasattr(cell, 'text_frame') and cell.text_frame:
                        self._process_text_frame(cell.text_frame, font_configs)
        except Exception as e:
            logger.error(f"处理表格时出错: {str(e)}")

    def _apply_font_config_to_run(self, run, font_configs: Dict):
        try:
            current_font = run.font
            current_latin_name = getattr(current_font, 'name', '') or ''
            current_size_pt = current_font.size.pt if current_font.size else 0

            current_ea_name = ""
            current_cs_name = ""

            try:
                if hasattr(run, '_r') and run._r is not None:
                    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                    rpr = run._r.find('.//a:rPr', ns)
                    if rpr is not None:
                        ea_elem = rpr.find('.//a:ea', ns)
                        current_ea_name = ea_elem.get('typeface', '') if ea_elem is not None else ''
                        cs_elem = rpr.find('.//a:cs', ns)
                        current_cs_name = cs_elem.get('typeface', '') if cs_elem is not None else ''
                        logger.debug(f"成功获取字体: EA='{current_ea_name}', CS='{current_cs_name}'")
                else:
                    logger.warning("Run对象没有_r属性，无法获取字体信息")
            except Exception as e:
                logger.error(f"获取XML字体信息失败: {str(e)}", exc_info=True)

            for config_name, config in font_configs.items():
                old_font = config.get('old_font')
                old_size = config.get('old_size')
                new_font = config.get('new_font')
                new_size = config.get('new_size')

                if new_font is None and new_size is None:
                    logger.warning(f"配置 '{config_name}' 无效：new_font 和 new_size 不能同时为空。")
                    continue

                font_match = False
                if old_font is not None:
                    if (current_latin_name == old_font or
                            current_ea_name == old_font or
                            current_cs_name == old_font):
                        font_match = True
                else:
                    font_match = True

                size_match = False
                if old_size is not None:
                    try:
                        if abs(float(old_size) - current_size_pt) < 0.1:
                            size_match = True
                    except (ValueError, TypeError):
                        size_match = False
                else:
                    size_match = True

                if not (font_match and size_match):
                    continue

                logger.info(f"文本 '{run.text[:20]}...' 匹配到配置 '{config_name}'。")

                if new_size is not None:
                    try:
                        new_size_val = float(new_size)
                        run.font.size = Pt(new_size_val)
                        logger.info(f"  -> 字号替换: {current_size_pt}pt -> {new_size_val}pt")
                    except (ValueError, TypeError) as e:
                        logger.error(f"  -> 设置新字号失败: {str(e)}")

                if new_font:
                    process_latin = config.get('latin', False)
                    process_ea = config.get('ea', False)
                    process_cs = config.get('cs', False)
                    logger.info(f"  -> 字体替换条件: latin={process_latin}, ea={process_ea}, cs={process_cs}")

                    rpr = None
                    if hasattr(run, '_r') and run._r is not None:
                        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                        rpr = run._r.find('.//a:rPr', ns)
                        if rpr is None:
                            rpr = ET.Element(f'{{{ns["a"]}}}rPr')
                            run._r.insert(0, rpr)
                        logger.info(f"  -> 找到 rPr 元素: {rpr}")
                    else:
                        logger.warning("Run对象没有有效的_r属性，无法处理字体配置")

                    if process_latin and (old_font is None or current_latin_name == old_font):
                        run.font.name = new_font
                        logger.info(f"  -> Latin 字体替换: '{current_latin_name}' -> '{new_font}'")

                    if process_ea and (old_font is None or current_ea_name == old_font):
                        if rpr is not None:
                            self._update_font_element(rpr, 'ea', new_font)
                            logger.info(f"  -> East-Asian 字体替换: '{current_ea_name}' -> '{new_font}'")
                        else:
                            logger.warning("rpr元素不存在，无法更新East-Asian字体")

                    if process_cs and (old_font is None or current_cs_name == old_font):
                        if rpr is not None:
                            self._update_font_element(rpr, 'cs', new_font)
                            logger.info(f"  -> Complex-Script 字体替换: '{current_cs_name}' -> '{new_font}'")
                        else:
                            logger.warning("rpr元素不存在，无法更新Complex-Script字体")

                    if run.text and len(run.text) > 0:
                        lang = 'zh-CN' if has_chinese(run.text) else 'en-US'
                        logger.info(f"  -> 检测到文本'{run.text}'语言: {lang}")
                        if hasattr(run, '_r') and run._r is not None:
                            xml_ns = run._r.nsmap.get('xml', 'http://www.w3.org/XML/1998/namespace')
                            lang_attr = '{' + xml_ns + '}lang'
                            run._r.attrib[lang_attr] = lang

                            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                            rpr = run._r.find('.//a:rPr', ns)
                            if rpr is not None:
                                rpr.set('lang', lang)
                                logger.info(f"  -> 已更新rPr语言设置: {lang}")
                            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                            p_element = run._r.getparent()
                            if p_element is not None:
                                endParaRPr = p_element.find('.//a:endParaRPr', ns)
                                if endParaRPr is not None and 'lang' in endParaRPr.attrib:
                                    endParaRPr.set('lang', lang)
                                    logger.info(f"  -> 已更新endParaRPr语言设置: {lang}")

                            logger.info(f"  -> '{run.text}'使用语言设置: {lang}")
                        else:
                            logger.warning("无法访问Run对象的XML元素，跳过语言设置")


        except Exception as e:
            logger.error(f"应用字体配置到 run 时发生未知错误: {str(e)}")

    def _update_font_element(self, rpr, font_type, new_font):
        if rpr is None:
            logger.error(f"更新{font_type}字体失败：rpr元素为None")
            return
        try:
            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            elem = rpr.find(f'.//a:{font_type}', ns)

            if elem is not None:
                elem.set('typeface', new_font)
            else:
                elem = ET.Element(f'{{{ns["a"]}}}{font_type}')
                elem.set('typeface', new_font)
                rpr.append(elem)
        except Exception as e:
            logger.error(f"更新{font_type}字体元素时出错: {str(e)}")

    def _process_gradient_effects(self, ppt_path: str, gradient_config: List[Dict], font_size: str = None,
                                  font_name: str = None):
        try:
            logging.info(f"开始处理渐变效果")

            temp_dir = self.file_manager.extract_pptx(ppt_path)
            logger.info(f"已解压到临时目录: {temp_dir}")

            slide_files = self.file_manager.get_slide_files(temp_dir)
            logging.info(f"找到 {len(slide_files)} 个幻灯片文件")

            for slide_file in slide_files:
                logging.info(f"处理幻灯片文件: {slide_file}")
                self._process_slide_gradient(slide_file, gradient_config, font_size, font_name)

            self.file_manager.compress_to_pptx(temp_dir, ppt_path)
            self.status_callback(f"渐变处理完成，已更新 {len(slide_files)} 个幻灯片")

        except Exception as e:
            logger.error(f"处理渐变效果时出错: {str(e)}")
            import traceback
            traceback.print_exc()
        finally:
            self.file_manager.cleanup()

    def _process_slide_gradient(self, slide_file: str, gradient_config: List[Dict] = None, font_size: str = None,
                                font_name: str = None):
        try:
            tree = self.xml_handler.load_xml(slide_file)
            if tree is None:
                logger.error(f"无法加载XML文件: {slide_file}")
                return

            target_configs = {}
            config_loaded = False

            config_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'config.json')
            try:
                with open(config_path, 'r', encoding='utf-8') as f:
                    config = json.load(f)
                    if isinstance(config, dict) and len(config) > 0:
                        target_configs = config
                        config_loaded = True
                        logging.info(f"从config.json加载了 {len(target_configs)} 个渐变方案")
            except FileNotFoundError:
                logger.warning("未找到config.json文件，将使用传入参数")
            except json.JSONDecodeError:
                logger.warning("config.json格式错误，将使用传入参数")
            except Exception as e:
                logger.error(f"加载config.json时出错: {str(e)}，将使用传入参数")

            if not config_loaded:
                if gradient_config and font_size and font_name:
                    target_configs = {font_size: {'gradient_config': gradient_config, 'font_name': font_name}}
                    logger.info(f"使用传入参数: 字号 {font_size} 字体 {font_name} 的渐变方案")
                else:
                    logger.warning("未提供有效的渐变配置，跳过处理")
                    return

            text_runs = self.xml_handler.find_text_runs(tree)
            logging.info(f"找到 {len(text_runs)} 个文本运行")

            gradient_applied_count = 0
            for text_run in text_runs:
                try:
                    rpr = text_run.find("a:rPr", self.xml_handler.namespaces)
                    if rpr is not None:
                        size = rpr.get('sz', '0')
                        if size != '0':
                            current_size = f"{int(size) / 100:g}"
                            logging.info(f"文本运行字号: {current_size}")

                            latin_font_elem = rpr.find("a:latin", self.xml_handler.namespaces)
                            ea_font_elem = rpr.find("a:ea", self.xml_handler.namespaces)
                            cs_font_elem = rpr.find("a:cs", self.xml_handler.namespaces)

                            current_fonts = []
                            if latin_font_elem is not None and latin_font_elem.get('typeface'):
                                current_fonts.append(latin_font_elem.get('typeface'))
                            if ea_font_elem is not None and ea_font_elem.get('typeface'):
                                current_fonts.append(ea_font_elem.get('typeface'))
                            if cs_font_elem is not None and cs_font_elem.get('typeface'):
                                current_fonts.append(cs_font_elem.get('typeface'))

                            unique_current_fonts = list(set(current_fonts))
                            logging.info(
                                f"文本运行字体: {', '.join(unique_current_fonts) if unique_current_fonts else 'None'}")

                            if current_size in target_configs:
                                config_entry = target_configs[current_size]
                                current_gradient = config_entry['gradient_config']
                                target_font = config_entry.get('font_name')

                                font_matched = any(font == target_font for font in unique_current_fonts)

                                if font_matched:
                                    self.xml_handler.apply_gradient_to_text_run(
                                        text_run, current_gradient, current_size, target_font)

                                    gradient_applied_count += 1
                                    logging.info(f"已应用渐变到字号 {current_size}、字体 {target_font} 的文本")
                                else:
                                    logging.warning(
                                        f"字体不匹配: 当前字体 {unique_current_fonts}, 目标字体 {target_font}")
                            else:
                                logging.warning(f"警告: 字号 {current_size} 不在渐变配置中")
                except Exception as e:
                    logger.error(f"处理文本运行时出错: {str(e)}")
                    continue

            self.status_callback(f"本幻灯片应用渐变的文本运行数量: {gradient_applied_count}")
            self.xml_handler.save_xml(tree, slide_file)
        except Exception as e:
            logger.error(f"处理渐变效果时出错: {str(e)}")
            import traceback
            traceback.print_exc()


def get_font_info_from_slide(self, slide_file: str) -> List[Dict]:
    try:
        prs = Presentation(slide_file)
        font_info = []

        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text.strip():
                                font = run.font
                                font_name = getattr(font, 'name', '') or ''
                                font_size = 0
                                if font.size:
                                    font_size = font.size.pt
                                ea_font = ''
                                cs_font = ''

                                try:
                                    if hasattr(run, '_element') and run._element is not None:
                                        rpr = run._element.find(
                                            './/{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
                                        if rpr is not None:
                                            ea_elem = rpr.find(
                                                './/{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
                                            cs_elem = rpr.find(
                                                './/{http://schemas.openxmlformats.org/drawingml/2006/main}cs')

                                            if ea_elem is not None:
                                                ea_font = ea_elem.get('typeface', '')
                                            if cs_elem is not None:
                                                cs_font = cs_elem.get('typeface', '')
                                except Exception as e:
                                    logging.error(f"获取详细字体信息时出错: {str(e)}")

                                font_info.append({
                                    'slide_idx': slide_idx + 1,
                                    'text': run.text[:50],
                                    'size': font_size,
                                    'latin_font': font_name,
                                    'ea_font': ea_font,
                                    'cs_font': cs_font
                                })

        return font_info

    except Exception as e:
        logging.error(f"提取字体信息时出错: {str(e)}")
        return []


def debug_font_info(self, input_path: str):
    font_info = self.get_font_info_from_slide(input_path)

    logging.info("\n=== PPT字体信息 ===")
    for info in font_info:
        logging.info(f"幻灯片 {info['slide_idx']}: '{info['text'][:20]}...' "
                     f"字号={info['size']} Latin='{info['latin_font']}' "
                     f"EA='{info['ea_font']}' CS='{info['cs_font']}'")
    logging.info("==================\n")
